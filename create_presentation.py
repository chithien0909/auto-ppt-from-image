#!/usr/bin/env python3

import os
import glob
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import time

def natural_sort_key(filename):
    """
    Generate a key for natural sorting of filenames with numbers.
    This ensures '2.png' comes before '10.png'
    """
    def convert(text):
        return int(text) if text.isdigit() else text.lower()
    
    return [convert(c) for c in re.split('([0-9]+)', filename)]

def show_progress(current, total, prefix='', suffix='', length=50, fill='█'):
    """
    Display a progress bar in the console.
    
    Args:
        current (int): Current progress value
        total (int): Total value
        prefix (str): Prefix string
        suffix (str): Suffix string
        length (int): Length of the progress bar
        fill (str): Character to fill the progress bar
    """
    percent = f"{100 * (current / float(total)):.1f}"
    filled_length = int(length * current // total)
    bar = fill * filled_length + '-' * (length - filled_length)
    print(f'\r{prefix} |{bar}| {percent}% {suffix}', end='\r')
    if current == total:
        print()

def create_presentation_from_images(images_folder="images", output_file="presentation.pptx", slide_format="4:3", fit_mode="contain"):
    """
    Create a PowerPoint presentation from images in a folder.
    
    Args:
        images_folder (str): Path to the folder containing images
        output_file (str): Name of the output PowerPoint file
        slide_format (str): Slide format, either "4:3" or "16:9"
        fit_mode (str): How to fit images on slides:
            - "cover": Image covers entire slide (may crop)
            - "contain": Entire image fits within slide (may have margins)
            - "stretch": Image stretches to fill slide (may distort)
    """
    start_time = time.time()
    
    # Create a new presentation with the specified aspect ratio
    if slide_format == "4:3":
        prs = Presentation()  # Default is 4:3
        # Explicitly set to 4:3 dimensions (10" x 7.5")
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        print(f"Creating presentation with 4:3 format ({prs.slide_width/Inches(1)}\" x {prs.slide_height/Inches(1)}\")")
    else:  # 16:9
        prs = Presentation()
        # Set to 16:9 dimensions (13.33" x 7.5")
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        print(f"Creating presentation with 16:9 format ({prs.slide_width/Inches(1)}\" x {prs.slide_height/Inches(1)}\")")
    
    image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.gif', '*.bmp', '*.tiff']
    image_files = []
    
    print("\nScanning for images...")
    for extension in image_extensions:
        image_files.extend(glob.glob(os.path.join(images_folder, extension)))
        image_files.extend(glob.glob(os.path.join(images_folder, extension.upper())))
    
    image_files.sort(key=lambda x: natural_sort_key(os.path.basename(x)))
    
    if not image_files:
        print(f"No image files found in {images_folder} folder.")
        return
    
    total_images = len(image_files)
    print(f"\nFound {total_images} image(s). Creating presentation in {slide_format} format with '{fit_mode}' fitting...")
    print("\nProcessing order:")
    for i, img_file in enumerate(image_files):
        print(f"  {i+1}. {os.path.basename(img_file)}")
    print()
    
    print("Creating slides...")
    for i, image_file in enumerate(image_files):
        show_progress(i + 1, total_images, prefix='Progress:', suffix=f'Slide {i+1}/{total_images}')
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        try:
            with Image.open(image_file) as img:
                img_width, img_height = img.size
            
            # Calculate aspect ratios
            slide_ratio = slide_width / slide_height
            img_ratio = img_width / img_height
            
            if fit_mode == "contain":
                # Fit entire image within slide (no cropping)
                if img_ratio > slide_ratio:
                    # Image is wider, scale by width
                    scale_factor = slide_width / img_width
                    final_width = slide_width
                    final_height = int(img_height * scale_factor)
                    left = 0
                    top = (slide_height - final_height) // 2
                else:
                    # Image is taller, scale by height
                    scale_factor = slide_height / img_height
                    final_height = slide_height
                    final_width = int(img_width * scale_factor)
                    left = (slide_width - final_width) // 2
                    top = 0
            
            elif fit_mode == "cover":
                # Cover entire slide (may crop image)
                if img_ratio > slide_ratio:
                    # Image is wider, scale by height
                    scale_factor = slide_height / img_height
                    final_height = slide_height
                    final_width = int(img_width * scale_factor)
                    left = (slide_width - final_width) // 2
                    top = 0
                else:
                    # Image is taller, scale by width
                    scale_factor = slide_width / img_width
                    final_width = slide_width
                    final_height = int(img_height * scale_factor)
                    left = 0
                    top = (slide_height - final_height) // 2
            
            elif fit_mode == "stretch":
                # Stretch to fill slide (may distort image)
                final_width = slide_width
                final_height = slide_height
                left = 0
                top = 0
            
            picture = slide.shapes.add_picture(
                image_file,
                left=left,
                top=top,
                width=final_width,
                height=final_height
            )
            
        except Exception as e:
            print(f"\nError adding image {image_file}: {str(e)}")
            continue
    
    print("\nSaving presentation...")
    try:
        prs.save(output_file)
        end_time = time.time()
        duration = end_time - start_time
        
        print(f"\nPresentation saved successfully as '{output_file}'")
        print(f"Created {len(prs.slides)} slides from {len(image_files)} images in {slide_format} format")
        print(f"Images fitted using '{fit_mode}' mode")
        print(f"Total processing time: {duration:.1f} seconds")
    except Exception as e:
        print(f"\nError saving presentation: {str(e)}")

def main():
    """Main function to run the script"""
    
    if not os.path.exists("images"):
        print("Error: 'images' folder not found in current directory.")
        return
    
    # Change fit_mode to "contain" to show entire image without cropping
    # Options: "contain" (fit entire image), "cover" (fill slide, may crop), "stretch" (fill slide, may distort)
    create_presentation_from_images(slide_format="4:3", fit_mode="contain")

if __name__ == "__main__":
    main() 